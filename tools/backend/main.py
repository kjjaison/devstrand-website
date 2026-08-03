"""
DevStrand Tools — PDF manipulation API
"""

from __future__ import annotations

import io
import os
import re
import shutil
import smtplib
import subprocess
import tempfile
import time
import zipfile
from collections import defaultdict, deque
from email.message import EmailMessage
from pathlib import Path
from typing import Literal
from uuid import uuid4

from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse, PlainTextResponse, Response
from fastapi.staticfiles import StaticFiles
from pypdf import PdfReader, PdfWriter

APP_ROOT = Path(__file__).resolve().parent.parent
FRONTEND = APP_ROOT / "frontend"
TMP = Path(os.environ.get("TOOLS_TMP", "/tmp/devstrand-tools"))
TMP.mkdir(parents=True, exist_ok=True)

MAX_UPLOAD_MB = int(os.environ.get("MAX_UPLOAD_MB", "100"))
MAX_COMPRESS_MB = int(os.environ.get("MAX_COMPRESS_MB", str(MAX_UPLOAD_MB)))
MAX_BYTES = MAX_UPLOAD_MB * 1024 * 1024
MAX_COMPRESS_BYTES = MAX_COMPRESS_MB * 1024 * 1024

# Email delivery (optional — disabled until SMTP_* is configured)
SMTP_HOST = os.environ.get("SMTP_HOST", "").strip()
SMTP_PORT = int(os.environ.get("SMTP_PORT", "587"))
SMTP_USER = os.environ.get("SMTP_USER", "").strip()
SMTP_PASSWORD = os.environ.get("SMTP_PASSWORD", "").strip()
SMTP_FROM = os.environ.get("SMTP_FROM", "").strip() or SMTP_USER
SMTP_TLS = os.environ.get("SMTP_TLS", "true").strip().lower() in {"1", "true", "yes"}
EMAIL_MAX_MB = int(os.environ.get("EMAIL_MAX_MB", "20"))
EMAIL_MAX_BYTES = EMAIL_MAX_MB * 1024 * 1024
EMAIL_RATE_LIMIT = int(os.environ.get("EMAIL_RATE_LIMIT", "8"))  # per IP per hour
EMAIL_RE = re.compile(r"^[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}$")
_email_hits: dict[str, deque[float]] = defaultdict(deque)

# Short-lived shareable download links
SHARES = TMP / "shares"
SHARES.mkdir(parents=True, exist_ok=True)
SHARE_ALLOWED_MINUTES = {15, 60, 360, 1440}  # 15m, 1h, 6h, 24h
SHARE_MAX_MB = int(os.environ.get("SHARE_MAX_MB", "50"))
SHARE_MAX_BYTES = SHARE_MAX_MB * 1024 * 1024
SHARE_RATE_LIMIT = int(os.environ.get("SHARE_RATE_LIMIT", "20"))
_share_hits: dict[str, deque[float]] = defaultdict(deque)

app = FastAPI(title="DevStrand Tools", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=[
        "Content-Disposition",
        "X-Original-Size",
        "X-Compressed-Size",
        "X-Compression-Level",
    ],
)


def _safe_name(name: str | None, fallback: str) -> str:
    base = Path(name or fallback).name
    return base.replace("..", "_")[:180] or fallback


def _check_size(upload: UploadFile, data: bytes, max_bytes: int | None = None, max_mb: int | None = None) -> None:
    limit = max_bytes if max_bytes is not None else MAX_BYTES
    label = max_mb if max_mb is not None else MAX_UPLOAD_MB
    if len(data) > limit:
        raise HTTPException(413, f"File too large. Max {label} MB.")


async def _read_upload(upload: UploadFile, max_bytes: int | None = None, max_mb: int | None = None) -> bytes:
    data = await upload.read()
    _check_size(upload, data, max_bytes=max_bytes, max_mb=max_mb)
    if not data:
        raise HTTPException(400, "Empty file.")
    return data


def _workdir() -> Path:
    path = TMP / uuid4().hex
    path.mkdir(parents=True, exist_ok=True)
    return path


def _libreoffice_convert(src: Path, out_dir: Path, target: str) -> Path:
    """Convert with LibreOffice headless. target e.g. pdf, docx."""
    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nolockcheck",
        "--nodefault",
        "--nofirststartwizard",
        "--convert-to",
        target,
        "--outdir",
        str(out_dir),
        str(src),
    ]
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
            env={**os.environ, "HOME": str(out_dir)},
        )
    except FileNotFoundError as exc:
        raise HTTPException(500, "LibreOffice is not installed in this container.") from exc
    except subprocess.TimeoutExpired as exc:
        raise HTTPException(504, "Conversion timed out.") from exc

    if result.returncode != 0:
        raise HTTPException(
            500,
            f"Conversion failed: {(result.stderr or result.stdout or 'unknown error')[:400]}",
        )

    outputs = list(out_dir.glob(f"{src.stem}.*"))
    outputs = [p for p in outputs if p.suffix.lower() != src.suffix.lower()]
    if not outputs:
        # LibreOffice sometimes keeps same stem with new ext
        candidates = [p for p in out_dir.iterdir() if p.is_file() and p != src]
        if not candidates:
            raise HTTPException(500, "Conversion produced no output.")
        return candidates[0]
    return outputs[0]


def _file_response(
    path: Path,
    download_name: str,
    media: str,
    extra_headers: dict[str, str] | None = None,
) -> FileResponse:
    return FileResponse(
        path,
        media_type=media,
        filename=download_name,
        background=None,
        headers=extra_headers or {},
    )


def _compress_with_ghostscript(src: Path, dst: Path, level: str) -> bool:
    """Return True if Ghostscript produced dst successfully."""
    if not shutil.which("gs"):
        return False

    # low = best quality / least shrink … maximum = smallest file
    pdfsettings = {
        "low": "/printer",
        "medium": "/ebook",
        "high": "/screen",
        "maximum": "/screen",
    }.get(level, "/ebook")

    cmd = [
        "gs",
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS={pdfsettings}",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        "-dDetectDuplicateImages=true",
        "-dCompressFonts=true",
        "-dSubsetFonts=true",
    ]
    if level == "maximum":
        cmd.extend(
            [
                "-dColorImageResolution=72",
                "-dGrayImageResolution=72",
                "-dMonoImageResolution=72",
                "-dDownsampleColorImages=true",
                "-dDownsampleGrayImages=true",
                "-dDownsampleMonoImages=true",
                "-dColorImageDownsampleType=/Bicubic",
                "-dGrayImageDownsampleType=/Bicubic",
            ]
        )
    elif level == "high":
        cmd.extend(
            [
                "-dColorImageResolution=100",
                "-dGrayImageResolution=100",
                "-dDownsampleColorImages=true",
                "-dDownsampleGrayImages=true",
            ]
        )

    cmd.extend([f"-sOutputFile={dst}", str(src)])
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180, check=False)
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False
    return result.returncode == 0 and dst.exists() and dst.stat().st_size > 0


def _compress_with_pypdf(data: bytes, level: str) -> bytes:
    """Fallback compressor: content-stream deflate + object dedupe."""
    zlib_level = {"low": 6, "medium": 9, "high": 9, "maximum": 9}.get(level, 9)
    writer = PdfWriter(clone_from=io.BytesIO(data))
    for page in writer.pages:
        try:
            page.compress_content_streams(level=zlib_level)
        except Exception:
            try:
                page.compress_content_streams()
            except Exception:
                pass
    try:
        writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)
    except Exception:
        pass
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _compress_pdf(data: bytes, level: str, work: Path) -> tuple[Path, str]:
    """
    Compress PDF. Prefers Ghostscript when installed; falls back to pypdf.
    Returns (output_path, engine_name).
    """
    level = (level or "medium").lower().strip()
    if level not in {"low", "medium", "high", "maximum"}:
        level = "medium"

    src = work / "input.pdf"
    src.write_bytes(data)
    out = work / "compressed.pdf"

    if _compress_with_ghostscript(src, out, level):
        # If GS somehow grew the file, keep the smaller of GS vs light pypdf.
        gs_bytes = out.read_bytes()
        if len(gs_bytes) >= len(data) and level in {"low", "medium"}:
            light = _compress_with_pypdf(data, level)
            if len(light) < len(gs_bytes):
                out.write_bytes(light)
                return out, "pypdf"
        return out, "ghostscript"

    compressed = _compress_with_pypdf(data, level)
    # Never return a larger file than the original for stream-only compress.
    if len(compressed) >= len(data):
        out.write_bytes(data)
    else:
        out.write_bytes(compressed)
    return out, "pypdf"


def _email_configured() -> bool:
    return bool(SMTP_HOST and SMTP_FROM)


def _client_ip(request: Request) -> str:
    forwarded = request.headers.get("cf-connecting-ip") or request.headers.get("x-forwarded-for")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.client.host if request.client else "unknown"


def _check_email_rate(ip: str) -> None:
    now = time.time()
    window = 3600.0
    hits = _email_hits[ip]
    while hits and now - hits[0] > window:
        hits.popleft()
    if len(hits) >= EMAIL_RATE_LIMIT:
        raise HTTPException(429, f"Too many emails from this address. Try again later (max {EMAIL_RATE_LIMIT}/hour).")
    hits.append(now)


def _check_share_rate(ip: str) -> None:
    now = time.time()
    window = 3600.0
    hits = _share_hits[ip]
    while hits and now - hits[0] > window:
        hits.popleft()
    if len(hits) >= SHARE_RATE_LIMIT:
        raise HTTPException(429, f"Too many share links. Try again later (max {SHARE_RATE_LIMIT}/hour).")
    hits.append(now)


def _purge_expired_shares() -> int:
    removed = 0
    now = time.time()
    if not SHARES.exists():
        return 0
    for folder in list(SHARES.iterdir()):
        if not folder.is_dir():
            continue
        meta_path = folder / "meta.json"
        expired = False
        if meta_path.exists():
            try:
                import json

                meta = json.loads(meta_path.read_text(encoding="utf-8"))
                if float(meta.get("expires_at", 0)) <= now:
                    expired = True
            except Exception:
                expired = True
        else:
            expired = True
        if expired:
            shutil.rmtree(folder, ignore_errors=True)
            removed += 1
    return removed


def _share_meta(token: str) -> dict | None:
    import json

    folder = SHARES / token
    meta_path = folder / "meta.json"
    data_path = folder / "file.bin"
    if not meta_path.exists() or not data_path.exists():
        return None
    try:
        meta = json.loads(meta_path.read_text(encoding="utf-8"))
    except Exception:
        return None
    if float(meta.get("expires_at", 0)) <= time.time():
        shutil.rmtree(folder, ignore_errors=True)
        return None
    meta["token"] = token
    meta["path"] = str(data_path)
    return meta


_purge_expired_shares()


def _send_attachment_email(*, to_email: str, filename: str, data: bytes, tool_name: str | None = None) -> None:
    if not _email_configured():
        raise HTTPException(
            503,
            "Email delivery is not configured. Set SMTP_HOST and SMTP_FROM on the tools server.",
        )
    if not EMAIL_RE.match(to_email):
        raise HTTPException(400, "Invalid email address.")
    if len(data) > EMAIL_MAX_BYTES:
        raise HTTPException(
            413,
            f"File is too large to email (max {EMAIL_MAX_MB} MB). Download instead, or compress first.",
        )

    safe_name = _safe_name(filename, "result.bin")
    subject_tool = (tool_name or "PDF tools").strip()[:80] or "PDF tools"
    msg = EmailMessage()
    msg["Subject"] = f"Your DevStrand file — {safe_name}"
    msg["From"] = SMTP_FROM
    msg["To"] = to_email
    msg.set_content(
        "Your file from DevStrand Tools is attached.\n\n"
        f"Tool: {subject_tool}\n"
        f"Filename: {safe_name}\n\n"
        "This message was sent because you asked to email the result on tools.devstrand.com.\n"
        "We do not keep a copy of your file after sending.\n\n"
        "— DevStrand (devstrand.com)\n"
    )
    maintype, subtype = "application", "octet-stream"
    lower = safe_name.lower()
    if lower.endswith(".pdf"):
        maintype, subtype = "application", "pdf"
    elif lower.endswith(".zip"):
        maintype, subtype = "application", "zip"
    elif lower.endswith(".docx"):
        maintype, subtype = "application", "vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif lower.endswith(".xlsx"):
        maintype, subtype = "application", "vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif lower.endswith(".pptx"):
        maintype, subtype = "application", "vnd.openxmlformats-officedocument.presentationml.presentation"
    elif lower.endswith((".jpg", ".jpeg")):
        maintype, subtype = "image", "jpeg"
    elif lower.endswith(".png"):
        maintype, subtype = "image", "png"
    msg.add_attachment(data, maintype=maintype, subtype=subtype, filename=safe_name)

    try:
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=45) as smtp:
            smtp.ehlo()
            if SMTP_TLS:
                smtp.starttls()
                smtp.ehlo()
            if SMTP_USER and SMTP_PASSWORD:
                smtp.login(SMTP_USER, SMTP_PASSWORD)
            smtp.send_message(msg)
    except smtplib.SMTPException as exc:
        raise HTTPException(502, f"Email send failed: {exc}") from exc
    except OSError as exc:
        raise HTTPException(502, f"Could not reach mail server: {exc}") from exc


@app.get("/api/health")
def health():
    has_lo = shutil.which("soffice") is not None
    has_gs = shutil.which("gs") is not None
    return {
        "ok": True,
        "libreoffice": has_lo,
        "ghostscript": has_gs,
        "max_upload_mb": MAX_UPLOAD_MB,
        "max_compress_mb": MAX_COMPRESS_MB,
        "email_enabled": _email_configured(),
        "email_max_mb": EMAIL_MAX_MB,
        "share_enabled": True,
        "share_max_mb": SHARE_MAX_MB,
        "share_ttl_minutes": sorted(SHARE_ALLOWED_MINUTES),
        "ocr_available": shutil.which("tesseract") is not None,
    }


@app.post("/api/email-result")
async def email_result(
    request: Request,
    file: UploadFile = File(...),
    to_email: str = Form(...),
    tool_name: str = Form(""),
):
    """Email a processed result file to the user. Not stored after send."""
    _check_email_rate(_client_ip(request))
    data = await file.read()
    if not data:
        raise HTTPException(400, "Empty file.")
    _send_attachment_email(
        to_email=to_email.strip(),
        filename=file.filename or "result.bin",
        data=data,
        tool_name=tool_name or None,
    )
    return {"ok": True, "message": f"Sent to {to_email.strip()}"}


@app.post("/api/share")
async def create_share(
    request: Request,
    file: UploadFile = File(...),
    expires_minutes: int = Form(60),
):
    """Store a result briefly and return a shareable download URL."""
    import json

    _purge_expired_shares()
    _check_share_rate(_client_ip(request))
    if expires_minutes not in SHARE_ALLOWED_MINUTES:
        raise HTTPException(
            400,
            f"expires_minutes must be one of: {', '.join(str(m) for m in sorted(SHARE_ALLOWED_MINUTES))}",
        )
    data = await file.read()
    if not data:
        raise HTTPException(400, "Empty file.")
    if len(data) > SHARE_MAX_BYTES:
        raise HTTPException(413, f"File too large to share (max {SHARE_MAX_MB} MB).")

    token = uuid4().hex
    folder = SHARES / token
    folder.mkdir(parents=True, exist_ok=True)
    (folder / "file.bin").write_bytes(data)
    filename = _safe_name(file.filename, "download.bin")
    content_type = file.content_type or "application/octet-stream"
    now = time.time()
    expires_at = now + expires_minutes * 60
    meta = {
        "filename": filename,
        "content_type": content_type,
        "size": len(data),
        "created_at": now,
        "expires_at": expires_at,
        "expires_minutes": expires_minutes,
    }
    (folder / "meta.json").write_text(json.dumps(meta), encoding="utf-8")
    return {
        "ok": True,
        "id": token,
        "url": f"/s/{token}",
        "download_url": f"/api/share/{token}/download",
        "expires_at": expires_at,
        "expires_in_seconds": expires_minutes * 60,
        "expires_minutes": expires_minutes,
        "filename": filename,
        "size": len(data),
    }


@app.get("/api/share/{token}")
def share_info(token: str):
    meta = _share_meta(token)
    if not meta:
        raise HTTPException(404, "Link not found or expired.")
    return {
        "ok": True,
        "filename": meta["filename"],
        "size": meta["size"],
        "content_type": meta["content_type"],
        "expires_at": meta["expires_at"],
        "expires_in_seconds": max(0, int(meta["expires_at"] - time.time())),
        "download_url": f"/api/share/{token}/download",
    }


@app.get("/api/share/{token}/download")
def share_download(token: str):
    meta = _share_meta(token)
    if not meta:
        raise HTTPException(404, "Link not found or expired.")
    return FileResponse(
        meta["path"],
        media_type=meta["content_type"],
        filename=meta["filename"],
    )


@app.get("/s/{token}", response_class=HTMLResponse)
def share_page(token: str):
    meta = _share_meta(token)
    if not meta:
        return HTMLResponse(
            """<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"/><meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Link expired | DevStrand Tools</title>
<link rel="stylesheet" href="/static/css/styles.css?v=13"/>
</head><body><div class="page-bg" aria-hidden="true"></div>
<main class="wrap share-page"><h1>This link has expired</h1>
<p class="lead">Shared files are deleted after a short time and are not kept permanently.</p>
<p><a class="btn" href="/">Back to tools</a></p></main></body></html>""",
            status_code=410,
        )
    remaining = max(0, int(meta["expires_at"] - time.time()))
    mins = remaining // 60
    secs = remaining % 60
    ttl = f"{mins}m {secs}s" if mins else f"{secs}s"
    name = meta["filename"]
    size_mb = meta["size"] / (1024 * 1024)
    size_label = f"{size_mb:.2f} MB" if size_mb >= 0.1 else f"{meta['size']} B"
    return HTMLResponse(
        f"""<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<meta name="robots" content="noindex,nofollow"/>
<title>Download {name} | DevStrand Tools</title>
<link rel="stylesheet" href="/static/css/styles.css?v=13"/>
</head><body><div class="page-bg" aria-hidden="true"></div>
<main class="wrap share-page">
<p class="eyebrow">Temporary share</p>
<h1>{name}</h1>
<p class="lead">{size_label} · expires in {ttl}</p>
<p class="privacy-note">This file is stored only until the link expires, then deleted. DevStrand does not keep permanent copies.</p>
<p class="share-actions"><a class="btn" href="/api/share/{token}/download">Download file</a>
<a class="btn btn-secondary" href="/">All tools</a></p>
</main></body></html>"""
    )


@app.post("/api/ocr")
async def ocr_pdf(file: UploadFile = File(...), language: str = Form("eng")):
    """Make a scanned PDF searchable with OCR (Tesseract via ocrmypdf)."""
    if shutil.which("tesseract") is None:
        raise HTTPException(503, "OCR is not available in this deployment (Tesseract missing).")
    try:
        import ocrmypdf
    except ImportError as exc:
        raise HTTPException(503, "OCR package is not installed.") from exc

    lang = re.sub(r"[^a-zA-Z_+-]", "", language or "eng")[:32] or "eng"
    data = await _read_upload(file)
    work = _workdir()
    src = work / "input.pdf"
    out = work / "ocr.pdf"
    src.write_bytes(data)
    try:
        ocrmypdf.ocr(
            src,
            out,
            language=lang,
            deskew=True,
            rotate_pages=True,
            force_ocr=False,
            skip_text=False,
            optimize=1,
            progress_bar=False,
        )
    except Exception as exc:
        # Retry forcing OCR for image-only / mixed docs
        try:
            ocrmypdf.ocr(
                src,
                out,
                language=lang,
                deskew=True,
                rotate_pages=True,
                force_ocr=True,
                optimize=1,
                progress_bar=False,
            )
        except Exception as exc2:
            raise HTTPException(400, f"OCR failed: {exc2}") from exc2
    if not out.exists():
        raise HTTPException(500, "OCR produced no output.")
    return _file_response(out, "ocr-searchable.pdf", "application/pdf")


@app.post("/api/esign")
async def esign_pdf(
    file: UploadFile = File(...),
    signature: UploadFile = File(...),
    signer_name: str = Form(""),
    pages: str = Form("all"),
    position: str = Form("bottom-right"),
    sig_width: float = Form(160),
):
    """Stamp a drawn/uploaded signature image onto PDF page(s)."""
    from PIL import Image as PILImage
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.colors import HexColor

    data = await _read_upload(file)
    sig_bytes = await signature.read()
    if not sig_bytes:
        raise HTTPException(400, "Signature image is required.")
    if len(sig_bytes) > 5 * 1024 * 1024:
        raise HTTPException(413, "Signature image too large (max 5 MB).")

    try:
        sig_img = PILImage.open(io.BytesIO(sig_bytes)).convert("RGBA")
    except Exception as exc:
        raise HTTPException(400, f"Invalid signature image: {exc}") from exc

    # Trim mostly-transparent edges from drawn signatures
    bbox = sig_img.getbbox()
    if bbox:
        sig_img = sig_img.crop(bbox)

    width_pt = max(40.0, min(float(sig_width), 400.0))
    aspect = sig_img.height / max(1, sig_img.width)
    height_pt = width_pt * aspect

    page_spec = (pages or "all").strip().lower()
    pos = (position or "bottom-right").strip().lower()
    if pos not in {"bottom-right", "bottom-left", "bottom-center", "center"}:
        raise HTTPException(400, "Invalid position.")

    work = _workdir()
    try:
        reader = PdfReader(io.BytesIO(data))
        writer = PdfWriter()
        total = len(reader.pages)
        target_indexes: set[int] = set()
        if page_spec in {"all", "*"}:
            target_indexes = set(range(total))
        else:
            try:
                for part in page_spec.split(","):
                    part = part.strip()
                    if not part:
                        continue
                    if "-" in part:
                        a, b = part.split("-", 1)
                        start, end = int(a), int(b)
                        for i in range(start, end + 1):
                            if 1 <= i <= total:
                                target_indexes.add(i - 1)
                    else:
                        i = int(part)
                        if 1 <= i <= total:
                            target_indexes.add(i - 1)
            except ValueError as exc:
                raise HTTPException(400, "Invalid pages value. Use all or e.g. 1,3,5-7.") from exc
        if not target_indexes:
            raise HTTPException(400, "No valid pages selected for signature.")

        sig_buf = io.BytesIO()
        sig_img.save(sig_buf, format="PNG")
        sig_buf.seek(0)
        sig_reader = ImageReader(sig_buf)

        for idx, page in enumerate(reader.pages):
            if idx in target_indexes:
                box = page.mediabox
                pw = float(box.width)
                ph = float(box.height)
                margin = 36.0
                if pos == "bottom-left":
                    x, y = margin, margin
                elif pos == "bottom-center":
                    x, y = (pw - width_pt) / 2, margin
                elif pos == "center":
                    x, y = (pw - width_pt) / 2, (ph - height_pt) / 2
                else:  # bottom-right
                    x, y = pw - width_pt - margin, margin

                packet = io.BytesIO()
                c = rl_canvas.Canvas(packet, pagesize=(pw, ph))
                c.drawImage(sig_reader, x, y + (18 if signer_name.strip() else 0), width=width_pt, height=height_pt, mask="auto")
                if signer_name.strip():
                    c.setFillColor(HexColor("#222222"))
                    c.setFont("Helvetica", 9)
                    c.drawString(x, y, signer_name.strip()[:80])
                c.save()
                packet.seek(0)
                overlay = PdfReader(packet).pages[0]
                page.merge_page(overlay)
            writer.add_page(page)

        out = work / "signed.pdf"
        with out.open("wb") as fh:
            writer.write(fh)
        return _file_response(out, "signed.pdf", "application/pdf")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"E-sign failed: {exc}") from exc


@app.post("/api/merge")
async def merge_pdfs(files: list[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(400, "Upload at least two PDF files.")
    writer = PdfWriter()
    work = _workdir()
    try:
        for f in files:
            data = await _read_upload(f)
            reader = PdfReader(io.BytesIO(data))
            for page in reader.pages:
                writer.add_page(page)
        out = work / "merged.pdf"
        with out.open("wb") as fh:
            writer.write(fh)
        return _file_response(out, "merged.pdf", "application/pdf")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"Merge failed: {exc}") from exc


@app.post("/api/split")
async def split_pdf(
    file: UploadFile = File(...),
    mode: Literal["pages", "ranges"] = Form("pages"),
    ranges: str = Form(""),
):
    """
    mode=pages → one PDF per page (zip)
    mode=ranges → ranges like 1-3,5,7-9 → zip of parts
    """
    data = await _read_upload(file)
    work = _workdir()
    try:
        reader = PdfReader(io.BytesIO(data))
        n = len(reader.pages)
        if n == 0:
            raise HTTPException(400, "PDF has no pages.")

        zip_path = work / "split.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            if mode == "pages":
                for i in range(n):
                    w = PdfWriter()
                    w.add_page(reader.pages[i])
                    buf = io.BytesIO()
                    w.write(buf)
                    zf.writestr(f"page-{i + 1}.pdf", buf.getvalue())
            else:
                parts = _parse_ranges(ranges, n)
                for idx, page_indexes in enumerate(parts, start=1):
                    w = PdfWriter()
                    for pi in page_indexes:
                        w.add_page(reader.pages[pi])
                    buf = io.BytesIO()
                    w.write(buf)
                    label = f"{page_indexes[0] + 1}-{page_indexes[-1] + 1}"
                    zf.writestr(f"part-{idx}_{label}.pdf", buf.getvalue())

        return _file_response(zip_path, "split.zip", "application/zip")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"Split failed: {exc}") from exc


def _parse_ranges(spec: str, page_count: int) -> list[list[int]]:
    if not spec.strip():
        raise HTTPException(400, "Provide ranges like 1-3,5,7-9")
    parts: list[list[int]] = []
    for chunk in spec.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            a, b = chunk.split("-", 1)
            start, end = int(a), int(b)
        else:
            start = end = int(chunk)
        if start < 1 or end > page_count or start > end:
            raise HTTPException(400, f"Invalid range {chunk} for {page_count} pages.")
        parts.append(list(range(start - 1, end)))
    if not parts:
        raise HTTPException(400, "No valid ranges.")
    return parts


@app.post("/api/compress")
async def compress_pdf(
    file: UploadFile = File(...),
    level: str = Form("medium"),
):
    data = await _read_upload(file, max_bytes=MAX_COMPRESS_BYTES, max_mb=MAX_COMPRESS_MB)
    work = _workdir()
    try:
        out, engine = _compress_pdf(data, level, work)
        original = len(data)
        compressed = out.stat().st_size
        return _file_response(
            out,
            "compressed.pdf",
            "application/pdf",
            extra_headers={
                "X-Original-Size": str(original),
                "X-Compressed-Size": str(compressed),
                "X-Compression-Level": (level or "medium").lower().strip(),
                "X-Compression-Engine": engine,
            },
        )
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"Compress failed: {exc}") from exc


@app.post("/api/watermark")
async def watermark_pdf(
    file: UploadFile = File(...),
    text: str = Form("DEVSTRAND"),
    opacity: float = Form(0.25),
):
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.colors import Color

    data = await _read_upload(file)
    work = _workdir()
    try:
        reader = PdfReader(io.BytesIO(data))
        writer = PdfWriter()
        for page in reader.pages:
            box = page.mediabox
            width = float(box.width)
            height = float(box.height)
            packet = io.BytesIO()
            c = rl_canvas.Canvas(packet, pagesize=(width, height))
            c.setFillColor(Color(0.18, 0.2, 0.45, alpha=max(0.05, min(opacity, 0.6))))
            c.saveState()
            c.translate(width / 2, height / 2)
            c.rotate(45)
            c.setFont("Helvetica-Bold", max(24, min(width, height) / 12))
            c.drawCentredString(0, 0, text[:80])
            c.restoreState()
            c.save()
            packet.seek(0)
            watermark = PdfReader(packet).pages[0]
            page.merge_page(watermark)
            writer.add_page(page)
        out = work / "watermarked.pdf"
        with out.open("wb") as fh:
            writer.write(fh)
        return _file_response(out, "watermarked.pdf", "application/pdf")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"Watermark failed: {exc}") from exc


@app.post("/api/edit")
async def edit_pdf(
    file: UploadFile = File(...),
    header_text: str = Form(""),
    footer_text: str = Form(""),
    rotate: int = Form(0),
):
    """Light edit: optional header/footer text and page rotation (0/90/180/270)."""
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.colors import HexColor

    data = await _read_upload(file)
    if rotate not in (0, 90, 180, 270):
        raise HTTPException(400, "rotate must be 0, 90, 180, or 270.")
    work = _workdir()
    try:
        reader = PdfReader(io.BytesIO(data))
        writer = PdfWriter()
        for page in reader.pages:
            if rotate:
                page.rotate(rotate)
            if header_text or footer_text:
                box = page.mediabox
                width = float(box.width)
                height = float(box.height)
                packet = io.BytesIO()
                c = rl_canvas.Canvas(packet, pagesize=(width, height))
                c.setFillColor(HexColor("#2e7dff"))
                c.setFont("Helvetica", 10)
                if header_text:
                    c.drawString(36, height - 28, header_text[:120])
                if footer_text:
                    c.drawString(36, 20, footer_text[:120])
                c.save()
                packet.seek(0)
                overlay = PdfReader(packet).pages[0]
                page.merge_page(overlay)
            writer.add_page(page)
        out = work / "edited.pdf"
        with out.open("wb") as fh:
            writer.write(fh)
        return _file_response(out, "edited.pdf", "application/pdf")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"Edit failed: {exc}") from exc


@app.post("/api/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...), dpi: int = Form(150)):
    from pdf2image import convert_from_bytes

    data = await _read_upload(file)
    dpi = max(72, min(dpi, 200))
    work = _workdir()
    try:
        images = convert_from_bytes(data, dpi=dpi, fmt="jpeg")
        if not images:
            raise HTTPException(400, "No pages rendered.")
        zip_path = work / "pages.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, img in enumerate(images, start=1):
                buf = io.BytesIO()
                img.convert("RGB").save(buf, format="JPEG", quality=85)
                zf.writestr(f"page-{i}.jpg", buf.getvalue())
        return _file_response(zip_path, "pdf-pages.zip", "application/zip")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"PDF to JPG failed: {exc}") from exc


@app.post("/api/jpg-to-pdf")
async def jpg_to_pdf(files: list[UploadFile] = File(...)):
    from PIL import Image

    if not files:
        raise HTTPException(400, "Upload one or more images.")
    work = _workdir()
    try:
        images = []
        for f in files:
            data = await _read_upload(f)
            img = Image.open(io.BytesIO(data)).convert("RGB")
            images.append(img)
        out = work / "images.pdf"
        first, rest = images[0], images[1:]
        first.save(out, save_all=True, append_images=rest)
        return _file_response(out, "images.pdf", "application/pdf")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"JPG to PDF failed: {exc}") from exc


@app.post("/api/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    data = await _read_upload(file)
    work = _workdir()
    src = work / _safe_name(file.filename, "input.pdf")
    src.write_bytes(data)

    # Prefer LibreOffice when available; fallback to pdf2docx
    try:
        if shutil.which("soffice"):
            out = _libreoffice_convert(src, work, "docx")
            return _file_response(out, f"{src.stem}.docx",
                                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    except HTTPException:
        pass

    try:
        from pdf2docx import Converter

        out = work / f"{src.stem}.docx"
        cv = Converter(str(src))
        cv.convert(str(out))
        cv.close()
        return _file_response(
            out,
            out.name,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    except Exception as exc:
        raise HTTPException(400, f"PDF to Word failed: {exc}") from exc


@app.post("/api/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    data = await _read_upload(file)
    work = _workdir()
    name = _safe_name(file.filename, "document.docx")
    if not name.lower().endswith((".doc", ".docx", ".odt", ".rtf")):
        raise HTTPException(400, "Upload a Word document (.doc, .docx).")
    src = work / name
    src.write_bytes(data)
    out = _libreoffice_convert(src, work, "pdf")
    return _file_response(out, f"{src.stem}.pdf", "application/pdf")


@app.post("/api/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    data = await _read_upload(file)
    work = _workdir()
    name = _safe_name(file.filename, "sheet.xlsx")
    if not name.lower().endswith((".xls", ".xlsx", ".ods", ".csv")):
        raise HTTPException(400, "Upload an Excel file (.xls, .xlsx).")
    src = work / name
    src.write_bytes(data)
    out = _libreoffice_convert(src, work, "pdf")
    return _file_response(out, f"{src.stem}.pdf", "application/pdf")


@app.post("/api/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Best-effort table extract via LibreOffice calc export when possible."""
    data = await _read_upload(file)
    work = _workdir()
    src = work / _safe_name(file.filename, "input.pdf")
    src.write_bytes(data)
    # LibreOffice can export PDF → HTML/writer; for spreadsheet use csv via pdfplumber fallback
    try:
        import pdfplumber
        from openpyxl import Workbook

        wb = Workbook()
        wb.remove(wb.active)
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                ws = wb.create_sheet(f"Page {i}")
                tables = page.extract_tables() or []
                row_n = 1
                if not tables:
                    text = page.extract_text() or ""
                    for line in text.splitlines():
                        ws.cell(row=row_n, column=1, value=line)
                        row_n += 1
                else:
                    for table in tables:
                        for row in table:
                            for col, cell in enumerate(row or [], start=1):
                                ws.cell(row=row_n, column=col, value=cell)
                            row_n += 1
                        row_n += 1
        out = work / f"{src.stem}.xlsx"
        wb.save(out)
        return _file_response(
            out,
            out.name,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
    except Exception as exc:
        raise HTTPException(400, f"PDF to Excel failed: {exc}") from exc


@app.post("/api/pdf-to-powerpoint")
async def pdf_to_powerpoint(file: UploadFile = File(...)):
    """Convert each PDF page to an image slide (PPTX)."""
    from pdf2image import convert_from_bytes
    from pptx import Presentation
    from pptx.util import Emu

    data = await _read_upload(file)
    work = _workdir()
    try:
        images = convert_from_bytes(data, dpi=120, fmt="png")
        if not images:
            raise HTTPException(400, "No pages rendered.")
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for img in images:
            slide = prs.slides.add_slide(blank)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            # Fit to slide
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            slide.shapes.add_picture(buf, Emu(0), Emu(0), width=slide_w, height=slide_h)
        out = work / "slides.pptx"
        prs.save(out)
        return _file_response(
            out,
            "slides.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"PDF to PowerPoint failed: {exc}") from exc


# Static frontend
if FRONTEND.exists():
    app.mount("/static", StaticFiles(directory=FRONTEND), name="static")

    @app.get("/")
    def index():
        html = (FRONTEND / "index.html").read_text(encoding="utf-8")
        return HTMLResponse(html)

    @app.get("/robots.txt", response_class=PlainTextResponse)
    def robots_txt():
        path = FRONTEND / "robots.txt"
        return path.read_text(encoding="utf-8") if path.exists() else "User-agent: *\nAllow: /\n"

    @app.get("/sitemap.xml")
    def sitemap_xml():
        path = FRONTEND / "sitemap.xml"
        body = path.read_text(encoding="utf-8") if path.exists() else (
            '<?xml version="1.0" encoding="UTF-8"?>\n'
            '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">'
            "<url><loc>https://tools.devstrand.com/</loc></url></urlset>"
        )
        return Response(content=body, media_type="application/xml")
